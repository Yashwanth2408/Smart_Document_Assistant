import os
import dotenv
import torch 
import fitz  
import pandas as pd
import docx
import shutil
from pptx import Presentation
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_chroma import Chroma
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import ChatGoogleGenerativeAI
import chromadb
from chromadb.config import Settings

# Load environment variables
dotenv.load_dotenv()
torch.set_num_threads(1)
os.environ["HF_HUB_DISABLE_SYMLINKS_WARNING"] = "1"

# API keys and configuration
GOOGLE_API_KEY = os.getenv("GOOGLE_API_KEY")
GEMINI_MODEL = "gemini-1.5-flash"
HUGGINGFACEHUB_API_TOKEN = os.getenv("HUGGINGFACEHUB_API_TOKEN")

if not GOOGLE_API_KEY:
    raise ValueError("❌ GOOGLE_API_KEY is missing! Please check .env file.")

# Initialize LLM
llm = ChatGoogleGenerativeAI(
    model=GEMINI_MODEL, 
    temperature=0.7, 
    google_api_key=GOOGLE_API_KEY  
)

# Initialize Embeddings (Hugging Face)
embedding_model = HuggingFaceEmbeddings(model_name="sentence-transformers/all-MiniLM-L6-v2")

# Initialize ChromaDB - Full reset approach to avoid persistence issues
vector_db_path = "chroma_db"

# Function to reset ChromaDB
def reset_chroma_db():
    global chroma_client, chroma_collection, vectorstore
    if 'chroma_client' in globals() and chroma_client:
        # Close the ChromaDB client
        chroma_client.close()
        print("ChromaDB client closed.")
    
    if os.path.exists(vector_db_path):
        try:
            shutil.rmtree(vector_db_path)
            print(f"Removed existing ChromaDB at {vector_db_path}")
        except Exception as e:
            print(f"Error deleting ChromaDB directory: {e}")
            # Try to delete the directory again after a short delay
            import time
            time.sleep(1)
            try:
                shutil.rmtree(vector_db_path)
                print(f"Removed existing ChromaDB at {vector_db_path}")
            except Exception as e:
                print(f"Failed to delete ChromaDB directory: {e}")
    
    # Re-initialize ChromaDB
    init_chroma_db()

# Create ChromaDB client with persistence settings
chroma_client = None
chroma_collection = None
vectorstore = None

def init_chroma_db():
    global chroma_client, chroma_collection, vectorstore
    if os.path.exists(vector_db_path):
        shutil.rmtree(vector_db_path)
        print(f"Removed existing ChromaDB at {vector_db_path}")
    
    chroma_client = chromadb.PersistentClient(
        path=vector_db_path,
        settings=Settings(
            anonymized_telemetry=False
        )
    )

    collection_name = "documents"
    try:
        chroma_collection = chroma_client.get_or_create_collection(name=collection_name)
        print(f"ChromaDB collection '{collection_name}' ready!")
    except Exception as e:
        print(f"Error initializing ChromaDB collection: {e}")
        # Create a new collection if error occurs
        if collection_name in [col.name for col in chroma_client.list_collections()]:
            chroma_client.delete_collection(name=collection_name)
        chroma_collection = chroma_client.create_collection(name=collection_name)
        print(f"Created new ChromaDB collection '{collection_name}'")

    vectorstore = Chroma(
        client=chroma_client,
        collection_name=collection_name,
        embedding_function=embedding_model
    )
    print("ChromaDB initialized successfully!")

# Initialize ChromaDB on first run
init_chroma_db()

# Functions to Extract Text from Different File Formats
def extract_text_from_pdf(pdf_path: str) -> str:
    text = ""
    try:
        with fitz.open(pdf_path) as doc:
            for page in doc:
                text += page.get_text("text") + "\n"
    except Exception as e:
        print(f"Error extracting text from PDF: {e}")
    return text

def extract_text_from_docx(docx_path: str) -> str:
    try:
        doc = docx.Document(docx_path)
        return "\n".join([para.text for para in doc.paragraphs])
    except Exception as e:
        print(f"Error extracting text from DOCX: {e}")
        return ""

def extract_text_from_txt(txt_path: str) -> str:
    try:
        with open(txt_path, "r", encoding="utf-8") as file:
            return file.read()
    except Exception as e:
        print(f"Error extracting text from TXT: {e}")
        return ""

def extract_text_from_excel(excel_path: str) -> str:
    try:
        df = pd.read_excel(excel_path)
        return "\n".join(df.astype(str).apply(lambda x: " | ".join(x), axis=1).tolist())
    except Exception as e:
        print(f"Error extracting text from Excel: {e}")
        return ""

def extract_text_from_pptx(pptx_path: str) -> str:
    try:
        prs = Presentation(pptx_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"Error extracting text from PPTX: {e}")
        return ""

# Process Uploaded File and Store in ChromaDB
def process_uploaded_file(file_path: str):
    global vectorstore
    ext = os.path.splitext(file_path)[-1].lower()
    
    # Extract text based on file extension
    if ext == ".pdf":
        text = extract_text_from_pdf(file_path)
    elif ext == ".docx":
        text = extract_text_from_docx(file_path)
    elif ext == ".txt":
        text = extract_text_from_txt(file_path)
    elif ext in [".xls", ".xlsx"]:
        text = extract_text_from_excel(file_path)
    elif ext == ".pptx":
        text = extract_text_from_pptx(file_path)
    else:
        print(f"❌ Unsupported file format: {ext}")
        return
    
    if not text.strip():
        print("❌ No text extracted from the document!")
        return

    print(f"📄 Extracted text from {file_path}:\n{text[:500]}...")  # Show first 500 characters

    # Store extracted text in ChromaDB
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=100)
    split_texts = text_splitter.split_text(text)
    
    if not split_texts:
        print("❌ No text chunks created!")
        return
    
    try:
        # Add to vectorstore directly
        vectorstore.add_texts(texts=split_texts)
        print(f"✅ Successfully added {len(split_texts)} text chunks to ChromaDB.")
    except Exception as e:
        print(f"❌ Error adding text to ChromaDB: {e}")

# Function to Search & Retrieve Answers
def search_documents(query):
    try:
        results = vectorstore.similarity_search(query, k=3)
        
        if not results:
            return "No relevant document content found."
        
        # Extracting text from Document objects
        return "\n\n".join([doc.page_content for doc in results])
    except Exception as e:
        print(f"❌ Error searching documents: {e}")
        return "Error searching the document database."

# Query Processing with Gemini
def answer_query_with_gemini(query):
    try:
        retrieved_text = search_documents(query)
        
        # If no relevant text is found, return a default response
        if retrieved_text == "No relevant document content found." or retrieved_text == "Error searching the document database.":
            return "I couldn't find relevant information in the document."

        # Generate AI Answer with Gemini
        prompt = f"""
        Based on the following document content, please answer the query.
        
        DOCUMENT CONTENT:
        {retrieved_text}
        
        QUERY:
        {query}
        
        Please provide a detailed and accurate response based only on the information in the document.
        """
        
        response = llm.invoke(prompt)
        return response.content
    except Exception as e:
        print(f"❌ Error generating answer: {e}")
        return f"Sorry, I encountered an error while trying to answer your question: {str(e)}"

# Main Script Execution (Only Runs When Executed Directly)
if __name__ == "__main__":
    file_path = input("\n📂 Enter the file path to upload: ").strip()
    if os.path.exists(file_path):
        process_uploaded_file(file_path)

        # Start Query Loop
        while True:
            user_input = input("\n🔍 Ask a question from the document (or type 'exit' to quit): ").strip()
            if user_input.lower() == "exit":
                print("👋 Goodbye!")
                break
            
            response = answer_query_with_gemini(user_input)
            print("\n🤖 AI:", response)
    else:
        print("❌ File not found. Please check the path and try again.")
